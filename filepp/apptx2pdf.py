from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

import os

def convert_pptx_to_pdf(input_file, output_file):
    # 创建一个 PDF 对象
    c = canvas.Canvas(output_file, pagesize=letter)

    # 打开 PPTX 文件
    prs = Presentation(input_file)

    for slide_num, slide in enumerate(prs.slides):
        image_file = f"slide_{slide_num}.png"
        slide.shapes._spTree.export(image_file, 'PNG')

        # 添加图像到 PDF
        c.drawImage(image_file, 0, 0, width=letter[0], height=letter[1])
        c.showPage()
        
        # 删除临时图像文件
        os.remove(image_file)

    # 保存 PDF 文件
    c.save()

if __name__ =="__main__":
    input_pptx = 'test/example.pptx'
    output_pdf = 'test/example.pdf'
    convert_pptx_to_pdf(input_pptx, output_pdf)

