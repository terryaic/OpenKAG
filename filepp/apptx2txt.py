import sys
from pptx import Presentation
from filepp.textutil import content_to_txt

def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    text_runs = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                tb = []
                for row in table.rows:
                    cells = [cell.text for cell in row.cells]
                    tb.append(cells)
                text_runs.append(("table",tb))
            elif hasattr(shape, "text"):
                text_runs.append(("paragraph",shape.text))
    return text_runs


def pptx_to_txt(file_path):
  content = extract_text_from_pptx(file_path)
  out = content_to_txt(content)
  return out

if __name__ == "__main__":
  # 使用示例
  file_path = sys.argv[1]
  out = pptx_to_txt(file_path)
  print(out)
